    
import sys
import os
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from transformers import TFAutoModelWithLMHead
from pptx import Presentation
from pptx.util import Inches
import re

from html.parser import HTMLParser
import urllib
import urllib.request
import requests
from selenium import webdriver
import time
import json
import codecs
from urllib.request import urlretrieve
import shutil
#pip install sentencepiece
#pip install protobuf
import xmind
from xmind.core.const import TOPIC_DETACHED
from xmind.core.markerref import MarkerId
from xmind.core.topic import TopicElement

title_page = 0
normal_page = 1
normal_page2 = 2
two_page = 3
interlude_page = 4
pic_page = 5

templates = {}
templates['ai'] = 1
templates['idea'] = 2
templates['abstract'] = 3
templates['color'] = 4

class Img_downloader(HTMLParser):

    def __init__(self):
        HTMLParser.__init__(self)
        self.cnt = 0

    @staticmethod
    def _attr(attrlist, attrname):
        for attr in attrlist:
            if attr[0] == attrname:
                return attr[1]
        return None

    def restart(self, ccnt):
        self.cnt = ccnt

    def handle_starttag(self, tag, attrs):
        if self.cnt > 0:
            global img_idx
            if tag == 'img' and 'mimg' in self._attr(attrs, 'class'):
                location = self._attr(attrs, 'src')
                urlretrieve(location, './image/{}.jpeg'.format(img_idx))
                img_idx += 1
                self.cnt -= 1


user_agent = 'Mozilla/4.0 (compatible; MSIE 5.5; Windows NT)'
headers = {'User-Agent': user_agent}
option = webdriver.ChromeOptions()
# 隐藏窗口
option.add_argument('headless')
# 防止打印一些无用的日志
option.add_experimental_option("excludeSwitches", ['enable-automation', 'enable-logging'])
driver = webdriver.Chrome("chromedriver.exe", chrome_options=option)
url = "https://www.bing.com/images/search?q="
img_downloader = Img_downloader()
generation_args = {}
img_idx = 0


def set_args():
    generation_args["num_beams"] = 3
    generation_args["early_stopping"] = True
    generation_args["num_return_sequences"] = 1
    generation_args["max_length"] = 20


def generate_title(context, model, tokenizer, args):
    inputs = tokenizer(text=context, return_tensors="pt")
    outputs = model.generate(**inputs, **args)
    return tokenizer.decode(outputs[0], skip_special_tokens=True)


def generate_abstract(context, model, tokenizer):
    inputs = tokenizer.encode(context, return_tensors="tf", max_length=512)
    outputs = model.generate(inputs, max_length=200, min_length=20, length_penalty=1.0, num_beams=4,
                             early_stopping=True)
    return tokenizer.decode(outputs[0], skip_special_tokens=True)


def get_data(t, content):
    passage = content.split('\n')
    print(passage)
    data = []
    # if t == '0':
    #     title = generate_title(content, title_model, title_tokenizer, generation_args)
    # else:
    #     title = t
    title = t
    for i in range(len(passage)):
        p = passage[i]

        abstract = generate_abstract(p, abstract_model, abstract_tokenizer)

        subtitle = generate_title(p, title_model, title_tokenizer, generation_args)

        data.append(subtitle + '.')
        data[i] += abstract
        data[i] = re.sub(r'\.|;|；', '\n', data[i])

    return title, data


def get_pic(title, ccnt):
    driver.get(url + title)
    res = driver.page_source
    img_downloader.restart(ccnt)
    img_downloader.feed(res)


def generate_ppt(author, title, data, style):
    style_val = templates[style]
    prs = Presentation('template' + str(style_val) + '.pptx')

    cover = prs.slides.add_slide(prs.slide_layouts[title_page])
    t = cover.shapes.title
    subtitle = cover.placeholders[1]
    t.text = title
    subtitle.text = author
    
    for i in range(len(data)):
        d = data[i].split('\n')
        if '' in d:
            d.remove('')

        interlude = prs.slides.add_slide(prs.slide_layouts[interlude_page])
        t = interlude.shapes.title
        t.text = d[0]

        if len(d) == 3:
            slide = prs.slides.add_slide(prs.slide_layouts[two_page])
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[normal_page])

        t = slide.shapes.title
        t.text = d[0]
        
        if len(d) == 3:
            get_pic(d[0], 1)
            img = './image/{}.jpeg'.format(img_idx - 1)
            pic = slide.shapes.add_picture(image_file=img, left=Inches(5.5), top=Inches(3.5),
                                        width=Inches(4), height=Inches(2))
            slide.shapes._spTree.insert(1, pic._element)

            for j in range(1, 3):
                content = slide.placeholders[j]
                content.text = d[j]

        else:
            get_pic(d[0], 1)
            img = './image/{}.jpeg'.format(img_idx - 1)
            pic = slide.shapes.add_picture(image_file=img, left=Inches(5.5), top=Inches(3.5),
                                        width=Inches(4), height=Inches(2))
            slide.shapes._spTree.insert(1, pic._element)

            content = slide.placeholders[1]
            for j in range(1, len(d)):
                content.text = content.text + '\n' + d[j]

    end = prs.slides.add_slide(prs.slide_layouts[interlude_page])
    t = end.shapes.title
    t.text = 'THANK YOU'

    filename = author + '.pptx'
    prs.save(filename)
    if os.path.exists('.\\static\\PPT\\' + filename):
        os.remove('.\\static\\PPT\\' + filename)
    shutil.move(filename, '.\\static\\PPT')


def generate_xmind(author, title, data):
    workbook = xmind.load('.\\static\\xmind\\' + author + '.xmind')

    first_sheet = workbook.getPrimarySheet()  # 获取第一个画布

    first_sheet.setTitle(title)   # 设置画布名称

    root_topic1 = first_sheet.getRootTopic()  # 获取画布中心主题，默认创建画布时会新建一个空白中心主题

    root_topic1.setTitle(title)   # 设置主题名称

    for i in range(len(data)):
        d = data[i].split('\n')
        if '' in d:
            d.remove('')

        sub_topic1 = root_topic1.addSubTopic() # 创建子主题，并设置名称
        sub_topic1.setTitle(d[0])

        for j in range(1, len(d)):
            sub_topic2 = sub_topic1.addSubTopic()
            sub_topic2.setTitle(d[j])

    xmind.save(workbook)


if __name__ == '__main__':
    with open('athletes.txt', encoding='utf-8',errors='ignore') as f:
        input_data = f.read()
        #print(input_data)
        tmp = input_data.split('\n\n')
        len = len(tmp)
        titles = []
        contents = []
        for i in range(1,len-1):
            titles.append(tmp[i].split('\t')[0])
            print(tmp[i].split('\t')[1])
            contents.append(tmp[i].split('\t')[1])
        
        
        set_args()
        # title_model_select = 'pbmstrk/t5-large-arxiv-abstract-title'
        # title_tokenizer = AutoTokenizer.from_pretrained(title_model_select)
        # title_model = AutoModelForSeq2SeqLM.from_pretrained(title_model_select)
        abstract_model_select = "t5-base"
        abstract_model = TFAutoModelWithLMHead.from_pretrained(abstract_model_select)
        abstract_tokenizer = AutoTokenizer.from_pretrained(abstract_model_select)
        for i in range(1,len-1):
            title, data = get_data(titles[i], contents[i])
    
            #title = 'abc'
            #data = ['apple\nred\nfresh', 'pineapple\nnotapple\nyellow']
            generate_ppt(titles[i], title[i], data, 1)
            generate_xmind(titles[i], title[i], data)

            with open('output.txt', 'w', encoding='utf-8') as fout:
                fout.write(title + '_' + name)
